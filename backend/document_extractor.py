# Copyright 2026 LLM Council Project
# Licensed under [LICENCE À DÉFINIR]
"""
Document extractor — router par type de fichier.
Détection via magic bytes (header) + extension de secours.
Moteurs chargés à la demande (lazy import).
"""

import os
from typing import Optional

# ── Magic bytes → type ────────────────────────────────────────────────────────
MAGIC = [
    (b"%PDF",                "pdf"),
    (b"PK\x03\x04",         "zip_based"),   # DOCX, XLSX, PPTX, ODT…
    (b"\xd0\xcf\x11\xe0",   "cfb"),         # DOC, XLS, PPT (ancien format)
    (b"\xef\xbb\xbf",       "utf8bom"),     # UTF-8 BOM → texte
    (b"\xff\xfe",            "utf16le"),     # UTF-16 LE → texte
    (b"\xfe\xff",            "utf16be"),     # UTF-16 BE → texte
]

def detect_type(path: str, filename: str) -> str:
    """Retourne le type détecté : 'pdf' | 'docx' | 'txt' | 'md' | 'unknown'."""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    # Lire les 8 premiers octets
    try:
        with open(path, "rb") as f:
            header = f.read(8)
    except Exception:
        return ext or "unknown"

    # PDF
    if header[:4] == b"%PDF":
        return "pdf"

    # ZIP-based (DOCX, PPTX, XLSX, ODT…) — distinguer par extension
    if header[:4] == b"PK\x03\x04":
        if ext in ("docx", "odt"):
            return "docx"
        if ext in ("pptx", "odp"):
            return "pptx"
        if ext in ("xlsx", "ods"):
            return "xlsx"
        return "docx"  # fallback raisonnable

    # Ancien format Microsoft (DOC, XLS, PPT)
    if header[:4] == b"\xd0\xcf\x11\xe0":
        return "doc_legacy"

    # Texte (UTF-8 BOM, UTF-16, ou extension texte explicite)
    if header[:3] in (b"\xef\xbb\xbf",) or header[:2] in (b"\xff\xfe", b"\xfe\xff"):
        return "txt"
    if ext in ("txt", "md", "markdown", "rst", "csv", "log", "json", "yaml", "yml", "toml", "ini", "xml", "html", "htm", "js", "ts", "py", "css"):
        return ext if ext in ("md", "csv") else "txt"

    return "unknown"


# ── Extracteurs ───────────────────────────────────────────────────────────────

def extract_pdf(path: str) -> str:
    import pdfplumber
    parts = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text = page.extract_text(x_tolerance=2, y_tolerance=2)
            if text:
                parts.append(text.strip())
            # Tableaux → représentation textuelle simple
            for table in page.extract_tables():
                rows = []
                for row in table:
                    rows.append(" | ".join(cell or "" for cell in row))
                if rows:
                    parts.append("\n".join(rows))
    return "\n\n".join(parts)


def extract_docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())
    # Tableaux
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            parts.append(" | ".join(cells))
    return "\n\n".join(parts)


def extract_doc_legacy(path: str) -> str:
    """Ancien .doc — tentative via LibreOffice si dispo, sinon erreur claire."""
    import subprocess, tempfile
    try:
        # Tenter une conversion via LibreOffice (souvent absent en dev)
        with tempfile.TemporaryDirectory() as tmpdir:
            subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "txt:Text", "--outdir", tmpdir, path],
                check=True, capture_output=True, timeout=30
            )
            base = os.path.splitext(os.path.basename(path))[0]
            out = os.path.join(tmpdir, base + ".txt")
            with open(out, "r", encoding="utf-8", errors="replace") as f:
                return f.read()
    except Exception:
        raise ValueError(
            "Format .doc (ancien) non supporté directement. "
            "Convertis en .docx dans Word ou LibreOffice avant d'uploader."
        )


def extract_txt(path: str, encoding: str = "utf-8") -> str:
    with open(path, "r", encoding=encoding, errors="replace") as f:
        return f.read()


def extract_md(path: str) -> str:
    return extract_txt(path)


def extract_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
    return "\n".join(texts)


def extract_xlsx(path: str) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    texts = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = "\t".join(str(c) for c in row if c is not None)
            if row_text.strip():
                texts.append(row_text)
    wb.close()
    return "\n".join(texts)


# ── Point d'entrée principal ──────────────────────────────────────────────────

def extract(path: str, filename: str, max_chars: int = 80_000) -> dict:
    """
    Extrait le contenu textuel d'un fichier.
    Retourne { content, type, truncated, char_count }.
    Lève ValueError avec un message clair si le type n'est pas supporté.
    """
    file_type = detect_type(path, filename)

    if file_type == "pdf":
        content = extract_pdf(path)
    elif file_type in ("docx",):
        content = extract_docx(path)
    elif file_type == "doc_legacy":
        content = extract_doc_legacy(path)
    elif file_type in ("txt", "utf8bom", "utf16le", "utf16be", "unknown"):
        content = extract_txt(path)
    elif file_type == "md":
        content = extract_md(path)
    elif file_type == "csv":
        content = extract_txt(path)
    elif file_type == "pptx":
        content = extract_pptx(path)
    elif file_type == "xlsx":
        content = extract_xlsx(path)
    else:
        raise ValueError(
            f"Type de fichier non supporté : {file_type} "
            f"(fichier : {filename}). "
            f"Types acceptés : PDF, DOCX, TXT, MD, CSV, PPTX, XLSX, XLS, ODS."
        )

    truncated = len(content) > max_chars
    if truncated:
        content = content[:max_chars] + f"\n\n[... contenu tronqué — {len(content):,} caractères au total, limite {max_chars:,}]"

    return {
        "content": content,
        "type": file_type,
        "truncated": truncated,
        "char_count": min(len(content), max_chars),
    }
